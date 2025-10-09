from pathlib import Path
from typing import List
from pptx import Presentation
from core.utils import Chunk, estimate_tokens

def parse_pptx(file_path: Path, max_tokens: int, overlap: int) -> List[Chunk]:
    chunks = []
    
    try:
        prs = Presentation(file_path)
        
        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_text_parts = []
            
            # Extract title
            if slide.shapes.title:
                slide_text_parts.append(f"Title: {slide.shapes.title.text}")
            
            # Extract text from shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_text_parts.append(shape.text)
            
            # Extract notes
            if slide.has_notes_slide:
                notes_text = slide.notes_slide.notes_text_frame.text
                if notes_text.strip():
                    slide_text_parts.append(f"Notes: {notes_text}")
            
            slide_text = '\n'.join(slide_text_parts)
            
            if slide_text.strip():
                chunks.append(Chunk(
                    text=slide_text,
                    section_path=f"slide_{slide_num}",
                    token_count=estimate_tokens(slide_text),
                    extra_meta={"slide": slide_num}
                ))
    except Exception as e:
        print(f"Error parsing PPTX {file_path}: {e}")
    
    return chunks